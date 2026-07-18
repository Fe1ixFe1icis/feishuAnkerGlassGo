"""Build the Anker GlassGo defense PPTX from outline and data."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "submission" / "assets"
DST = ROOT / "submission" / "Anker_GlassGo_答辩PPT.pptx"

DARK = RGBColor(10, 25, 47)
ACCENT = RGBColor(0, 150, 255)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(180, 190, 200)


def set_slide_bg(slide, color=DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=16, color=WHITE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return box


def make_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_textbox(slide, 0.5, 2.5, 9, 1.2, title, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 0.5, 4.0, 9, 0.8, subtitle, font_size=20, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


def make_content_slide(prs, title, bullets, note=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_textbox(slide, 0.5, 0.4, 9, 0.8, title, font_size=28, bold=True, color=WHITE)
    add_bullets(slide, 0.5, 1.4, 9, 5.5, bullets, font_size=18, color=WHITE)
    if note:
        add_textbox(slide, 0.5, 6.8, 9, 0.5, note, font_size=12, color=GRAY)
    return slide


def make_two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_textbox(slide, 0.5, 0.4, 9, 0.8, title, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, 0.5, 1.3, 4.2, 0.5, left_title, font_size=20, bold=True, color=ACCENT)
    add_bullets(slide, 0.5, 1.9, 4.2, 5.0, left_bullets, font_size=15, color=WHITE)
    add_textbox(slide, 5.3, 1.3, 4.2, 0.5, right_title, font_size=20, bold=True, color=ACCENT)
    add_bullets(slide, 5.3, 1.9, 4.2, 5.0, right_bullets, font_size=15, color=WHITE)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    make_title_slide(prs, "Anker GlassGo", "AI 眼镜的跨品牌补能系统 · AI 原生产品定义工作流")

    # 2. Elevator pitch
    pitch = (ASSETS / "elevator_pitch.md").read_text(encoding="utf-8")
    pitch_lines = [line.strip().lstrip("- ").replace("**", "") for line in pitch.splitlines() if line.strip() and not line.startswith("#")]
    make_content_slide(prs, "电梯演讲", pitch_lines[:6])

    # 3. Market storm
    make_content_slide(
        prs,
        "市场风暴：2025 年全球 AI 眼镜出货约 870 万台",
        [
            "同比增速超过 300%，Meta、雷鸟、Rokid、Xreal 等品牌百镜大战",
            "主机厂各自封闭补能方案：触点、仓体、协议互不兼容",
            "2026 年中国将智能眼镜纳入国补，进一步刺激换机潮",
            "配件生态滞后主机设备 1–2 年，正是跨品牌补能的窗口期",
        ],
        note="数据来源：Omdia / IDC / TrendForce / 洛图科技（公开链接待补充）",
    )

    # 4. Fracture
    make_content_slide(
        prs,
        "断裂带：用户真正的痛点不是续航，而是充电时失明",
        [
            "旅行 Vlogger：放进盒子充电 = 拍摄中断，景色不等人",
            "商务用户：会议翻译时低电，摘下眼镜等于断线",
            "近视用户：充电 2 小时，看黑板都是糊的",
            "市场上没有跨品牌、边用边充的能源确定性方案",
        ],
    )

    # 5. Product system
    make_content_slide(
        prs,
        "产品系统：Pocket + Clip + App",
        [
            "GlassGo Pocket：随身充电盒 + 可替换触点舌片，15 分钟快充 80%",
            "GlassGo Clip：<60g 夹持式边充电池，边拍摄/翻译边补电",
            "GlassGo App：电量看板 + 低电预警 + 舌片识别购买 + 投票开模",
            "定价 $59.99，首年目标 15–25 万套",
        ],
    )

    # 6. Workflow overview
    make_content_slide(
        prs,
        "AI 原生工作流：六阶段 × 七角色",
        [
            "S0 提问教练：把原始需求框定为可分析的产品问题",
            "S1 情报分析师：四路并行采集财报/评测/差评/政策信号",
            "S2 机会侦探：聚类信号并量化打分，选出最值得下注的假设",
            "S3 概念设计师：多形态发散，禁止首次收敛",
            "S4 用户替身红队：5 个数据驱动 persona 第一人称攻击",
            "S5 专家委员会：技术 / 供应链 / 合规 / 市场并行评审",
            "S6 书记官：综合所有材料与人类决策，输出结构化提案",
        ],
    )

    # 7. Memory isolation
    make_content_slide(
        prs,
        "记忆隔离 + 证据链：防止 AI 说服 AI",
        [
            "每个 Agent 只能看到上游阶段产出的结构化结果，看不到中间推理",
            "Fact 与 Inference 分离：事实必须附可核查信源，推理必须附置信度",
            "人类在 S0–S5 六个断点做最终价值判断，AI 不替人拍板",
            "所有决策与证据写入 state，可审计、可复盘、可迭代",
        ],
    )

    # 8. Red team
    make_content_slide(
        prs,
        "红队对抗实录：每条攻击都改写定义",
        [
            "张野（旅行 Vlogger）：放进 Pocket 充电 = 拍摄中断 → 诞生 GlassGo Clip",
            "Linda（商务用户）：触点不统一怎么办 → 可替换舌片 + 投票开模",
            "阿凯（数码尝鲜党）：眼镜都退了谁买配件 → 瞄准留存用户 ARPU",
            "王姐（教师/近视）：充电时看不见 → 15 分钟快充 80% 成核心指标",
            "陆工（硬件 PM）：Meta 二代续航提升 → 升级为跨品牌能源确定性叙事",
        ],
    )

    # 9. Key pivot
    make_content_slide(
        prs,
        "关键修正：从续航不足到能源确定性",
        [
            "初始假设：AI 眼镜续航短，用户需要更大电池",
            "红队攻击后发现：真实痛点是充电时必须摘下",
            "机会重新定义：做跨品牌、跨场景、边用边充的能源基础设施",
            "产品形态从大一统充电宝收敛为通用仓体 + 可替换舌片 + 边充 Clip",
        ],
    )

    # 10. Market sizing
    make_content_slide(
        prs,
        "市场测算：2027 年三情景",
        [
            "保守情景：渗透率 5%，营收约 1.1 亿美元",
            "基准情景：渗透率 8%，营收约 1.9 亿美元",
            "乐观情景：渗透率 12% + B2B 模组授权，营收约 2.8 亿美元",
            "首年通过众筹做止损闸，验证付费意愿与触点兼容性后再量产",
        ],
        note="测算为提案级假设，需在 MVP 阶段校准",
    )

    # 11. Roadmap
    make_content_slide(
        prs,
        "路线图：从众筹 MVP 到 B2B 标准",
        [
            "P0（0–6 个月）：Kickstarter MVP，5 个主流机型适配，验证付费意愿",
            "P1（6–18 个月）：量产 + 生态扩展，新增 10+ 机型，推出 GlassGo Tower",
            "P2（18–36 个月）：B2B 充电模组授权，成为中小眼镜品牌标配",
        ],
    )

    # 12. Methodology comparison
    make_two_col_slide(
        prs,
        "方法论对比：AI 驱动 vs 经验驱动",
        "传统 · 经验驱动",
        [
            "信息：工单、销售数据、焦点小组",
            "判断：产品经理直觉 + 高管经验",
            "生成：内部脑暴，路径依赖",
            "验证：评审会，资历权重",
            "周期：月级反复开会",
        ],
        "AI 原生 · 证据 + 对抗驱动",
        [
            "信息：四路并行自动聚类信号",
            "判断：量化打分，可审计复现",
            "生成：AI 多形态发散",
            "验证：用户替身第一人称攻击",
            "周期：单日跑通 S0–S6",
        ],
    )

    # 13. Why Anker
    make_content_slide(
        prs,
        "为什么适合安克：能力圆心 + 浅海战略 + 安全合规",
        [
            "能力圆心：安克是全球领先的充电品牌，2025 年充电储能类营收 154.02 亿元",
            "浅海战略：AI 眼镜配件是新兴蓝海，竞品尚未形成品牌心智",
            "安全合规：10Wh 以内电池符合航空携带标准，安克有成熟全球认证经验",
            "中立定位：不站队任何主机厂，做品类的能源基础设施",
        ],
    )

    # 14. Closing
    make_content_slide(
        prs,
        "谢谢聆听 · 代码与 Demo",
        [
            "完整代码：src/glassgo_workflow/（LangGraph + Pydantic + 6 个人类断点）",
            "运行命令：py scripts/run_cli.py --auto",
            "交互 Demo：index.html / workflow-demo.html",
            "测试：pytest tests/（6 项全部通过）",
        ],
    )

    prs.save(DST)
    print(f"Submission PPTX saved to: {DST}")


if __name__ == "__main__":
    main()
